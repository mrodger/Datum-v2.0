#!/usr/bin/env python3
"""Datum Codex — FastAPI server (bridge proxy backend).

1:1 functional clone of belle_assistant/belle-api/server.py with two changes:
  1. run_agent() receives full conversation history (stateless API needs it).
  2. Memory extraction + title generation use httpx → OpenAI API.
"""
import asyncio
import ipaddress
import json
import os
import re
import time
import uuid
from collections import defaultdict
from pathlib import Path
from datetime import datetime, timezone

import httpx

OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY', '')
AUX_MODEL = os.environ.get('AUX_MODEL', 'gpt-4o-mini')
from fastapi import Depends, FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from starlette.middleware.base import BaseHTTPMiddleware

from db import (
    delete_memory,
    get_all_memory,
    get_history,
    init_db,
    list_memory,
    list_recent_conversations,
    log_message,
    set_conversation_title,
    store_memory,
)

# ── App setup ─────────────────────────────────────────────────────────────────

app = FastAPI(
    title="belle-assistant-openrouter",
    docs_url=None,
    redoc_url=None,
    openapi_url=None,
)


class SecurityHeadersMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["Referrer-Policy"] = "no-referrer"
        response.headers["Permissions-Policy"] = "geolocation=(), camera=(), microphone=()"
        response.headers["Content-Security-Policy"] = (
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline' https://unpkg.com https://cdn.jsdelivr.net; "
            "style-src 'self' 'unsafe-inline' https://unpkg.com https://cdn.jsdelivr.net; "
            "font-src 'self' https://cdn.jsdelivr.net https://fonts.gstatic.com; "
            "img-src 'self' data: https:; "
            "connect-src 'self' https: http://localhost:8130 http://192.168.88.102:8130; "
            "frame-src 'self' blob: http://localhost:8130 http://192.168.88.102:8130;"
        )
        response.headers["Server"] = "datum-codex"
        return response


app.add_middleware(SecurityHeadersMiddleware)
init_db()

# ── Auth / rate limiting ───────────────────────────────────────────────────────

TRUSTED_IPS = {"127.0.0.1", "::1"}


def _require_trusted_ip(request: Request):
    ip = request.headers.get("X-Real-IP") or (
        request.client.host if request.client else ""
    )
    try:
        ip_obj = ipaddress.ip_address(ip)
        trusted = ip in TRUSTED_IPS or ip_obj.is_private or ip_obj.is_loopback
    except ValueError:
        trusted = False
    if not trusted:
        raise HTTPException(403, "Forbidden")


_rate_limits: dict = defaultdict(list)
RATE_LIMIT_RPM = 10
RATE_LIMIT_UPLOAD_RPM = 10


def _check_rate_limit(key: str, limit: int):
    now = time.time()
    _rate_limits[key] = [t for t in _rate_limits[key] if now - t < 60]
    if len(_rate_limits[key]) >= limit:
        raise HTTPException(429, "Rate limit exceeded. Try again shortly.")
    _rate_limits[key].append(now)


# ── Memory extraction via OpenAI API ─────────────────────────────────────────────

# Use a cheap fast model for fire-and-forget auxiliary calls
MEMORY_EXTRACT_MODEL = os.environ.get("AUX_MODEL", "gpt-4o-mini")
TITLE_MODEL = os.environ.get("AUX_MODEL", "gpt-4o-mini")

MEMORY_EXTRACT_PROMPT = """You are a memory extractor for a marketing assistant serving a Kauai boutique agency.

Given the following user question and assistant response, extract 0–3 facts worth remembering across sessions.
Only extract things that are specific, reusable, and not already obvious from the system prompt.
Examples: a client name or niche, a brand voice preference, a content strategy decision, a recurring campaign, a platform the client prioritises.

Return ONLY a JSON array. Empty array if nothing is worth storing. No explanation, no markdown.
Format: [{{"key": "short_key", "value": "concise fact", "category": "one of: client|brand_voice|strategy|platform|campaign|preference"}}]

User: {user_msg}
Assistant: {assistant_msg}"""

TITLE_PROMPT = """Generate a short, descriptive title (4-8 words) for this marketing assistant chat session based on the first exchange.
Be specific — mention the client, platform, campaign type, or content format if clear.
Return ONLY the title text, no quotes, no punctuation at the end.

User: {user_msg}
Assistant: {assistant_msg}"""


async def _litellm_complete(prompt: str, model: str, max_tokens: int = 500) -> str:
    """Non-streaming completion via OpenAI API for auxiliary tasks."""
    if not OPENAI_API_KEY:
        return ""
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.post(
                "https://api.openai.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
                json={
                    "model": model,
                    "messages": [{"role": "user", "content": prompt}],
                    "max_tokens": max_tokens,
                },
            )
            resp.raise_for_status()
            return resp.json()["choices"][0]["message"]["content"].strip()
    except Exception:
        return ""


async def extract_and_store_memories(user_msg: str, assistant_msg: str):
    """Fire-and-forget: extract memorable facts from this exchange."""
    if not user_msg or not assistant_msg or len(assistant_msg) < 50:
        return
    prompt = MEMORY_EXTRACT_PROMPT.format(
        user_msg=user_msg[:1000],
        assistant_msg=assistant_msg[:2000],
    )
    try:
        text = await _litellm_complete(prompt, MEMORY_EXTRACT_MODEL, max_tokens=500)
        if not text or len(text) > 10_000:
            return
        start = text.find("[")
        end = text.rfind("]") + 1
        if start == -1 or end == 0:
            return
        facts = json.loads(text[start:end])
        for f in facts:
            if f.get("key") and f.get("value"):
                store_memory(
                    str(f["key"])[:100],
                    str(f["value"])[:500],
                    str(f.get("category", "general"))[:50],
                )
    except Exception:
        pass


async def generate_title_if_new(conversation_id: str, user_msg: str, assistant_msg: str):
    """Fire-and-forget: name an untitled conversation."""
    if not user_msg or not assistant_msg:
        return
    convs = list_recent_conversations(limit=100)
    existing = next((c for c in convs if c["id"] == conversation_id), None)
    if existing and existing.get("title"):
        return
    prompt = TITLE_PROMPT.format(
        user_msg=user_msg[:500],
        assistant_msg=assistant_msg[:1000],
    )
    try:
        title = await _litellm_complete(prompt, TITLE_MODEL, max_tokens=30)
        title = title.strip().strip('"').strip("'")[:80]
        if title:
            set_conversation_title(conversation_id, title)
    except Exception:
        pass


# ── Static files ──────────────────────────────────────────────────────────────

static_dir = Path(__file__).parent / "static"
if static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

FILES_DIR = Path(os.environ.get("FILES_DIR", "/app/files"))
ALLOWED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".pdf", ".docx", ".xlsx", ".xls",
    ".png", ".jpg", ".jpeg", ".gif", ".svg", ".geojson", ".html", ".htm", ".pptx",
}
MAX_UPLOAD_MB = 50

ALLOWED_MODELS = {
    "gpt-5.4",
    "gpt-4o",
    "gpt-4o-mini",
    "gpt-4-turbo",
    "claude-sonnet-4-6",
    "claude-opus-4-6",
    "claude-opus-4-7",
    "claude-haiku-4-5",
    "default",
    "gpt-5.4-mini",
    "codex-5.4",
}

UUID_RE = re.compile(
    r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$",
    re.IGNORECASE,
)


# ── Request models ────────────────────────────────────────────────────────────

class ChatRequest(BaseModel):
    messages: list
    model: str = "default"
    conversation_id: str = None


class MemoryRequest(BaseModel):
    key: str = Field(..., min_length=1, max_length=100)
    value: str = Field(..., min_length=1, max_length=500)
    category: str = Field(default="general", max_length=50)


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/")
async def root():
    index = static_dir / "index.html"
    if index.exists():
        return FileResponse(
            str(index),
            media_type="text/html",
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return {"status": "ok", "service": "belle-assistant-openrouter"}


@app.get("/sw.js")
async def service_worker():
    return FileResponse(str(static_dir / "sw.js"), media_type="application/javascript")


@app.get("/manifest.json")
async def manifest():
    return FileResponse(str(static_dir / "manifest.json"), media_type="application/manifest+json")


@app.get("/icon-192.png")
async def icon192():
    return FileResponse(str(static_dir / "icon-192.png"), media_type="image/png")


@app.get("/icon-512.png")
async def icon512():
    return FileResponse(str(static_dir / "icon-512.png"), media_type="image/png")


@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "datum-codex",
        "backend": "codex-bridge",
        "model": os.environ.get("DEFAULT_MODEL", "gpt-5.4"),
    }


@app.post("/chat")
async def chat(req: ChatRequest, request: Request):
    from agent import run_agent

    client_ip = request.client.host if request.client else "unknown"
    _check_rate_limit(f"chat:{client_ip}", RATE_LIMIT_RPM)

    if not req.messages:
        raise HTTPException(400, "messages required")

    last = req.messages[-1]
    if not isinstance(last, dict) or "content" not in last or "role" not in last:
        raise HTTPException(400, "invalid message format")
    if last.get("content") and len(str(last["content"])) > 20_000:
        raise HTTPException(400, "message too long: max 20,000 characters")

    if req.conversation_id:
        if not UUID_RE.match(req.conversation_id):
            raise HTTPException(400, "invalid conversation_id format")
        conversation_id = req.conversation_id
    else:
        conversation_id = str(uuid.uuid4())

    # Log user message to DB before loading history
    if last.get("role") == "user" and last.get("content") and last["content"] != "init":
        log_message(conversation_id, "user", str(last["content"]))

    # Load full history (includes current message just logged)
    history = get_history(conversation_id)
    memory_context = get_all_memory()

    model = req.model if req.model in ALLOWED_MODELS else "default"

    async def stream():
        assistant_text = []
        try:
            async for event in run_agent(
                req.messages,
                model=model,
                conversation_id=conversation_id,
                memory_context=memory_context,
                history=history,
            ):
                if event["type"] == "text":
                    assistant_text.append(event["text"])
                yield f"data: {json.dumps(event)}\n\n"
        except Exception:
            yield f"data: {json.dumps({'type': 'error', 'detail': 'An internal error occurred. Please try again.'})}\n\n"
        finally:
            full_text = "".join(assistant_text)
            if full_text:
                log_message(conversation_id, "assistant", full_text)
                user_msg = (
                    str(last.get("content", "")) if last.get("role") == "user" else ""
                )
                asyncio.create_task(extract_and_store_memories(user_msg, full_text))
                asyncio.create_task(
                    generate_title_if_new(conversation_id, user_msg, full_text)
                )
            yield f'data: {{"type": "done", "conversation_id": "{conversation_id}"}}\n\n'

    return StreamingResponse(
        stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/upload")
async def upload(file: UploadFile = File(...), request: Request = None):
    if request:
        client_ip = request.client.host if request.client else "unknown"
        _check_rate_limit(f"upload:{client_ip}", RATE_LIMIT_UPLOAD_RPM)

    if not file.filename:
        raise HTTPException(400, "no filename")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"file type not allowed: {ext}")

    max_bytes = MAX_UPLOAD_MB * 1024 * 1024
    chunks = []
    total = 0
    while True:
        chunk = await file.read(64 * 1024)
        if not chunk:
            break
        total += len(chunk)
        if total > max_bytes:
            raise HTTPException(400, f"file too large: max {MAX_UPLOAD_MB} MB")
        chunks.append(chunk)
    content = b"".join(chunks)

    FILES_DIR.mkdir(parents=True, exist_ok=True)
    dest = FILES_DIR / Path(file.filename).name
    dest.write_bytes(content)

    return {"status": "ok", "filename": dest.name, "size_kb": len(content) // 1024}


@app.get("/read-file/{filename}")
def read_file_for_chat(filename: str):
    """Return extracted text of an uploaded file for chat injection."""
    safe = Path(filename).name
    path = FILES_DIR / safe
    if not path.exists():
        raise HTTPException(404, "file not found")
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
        elif ext in (".docx",):
            from docx import Document
            doc = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        else:
            text = path.read_text(errors="replace")
        return {"filename": safe, "content": text[:50_000]}
    except Exception as e:
        raise HTTPException(500, f"Failed to read file: {e}")


@app.get("/read-excel/{filename}")
def read_excel_structured(filename: str):
    """Return Excel data as structured JSON."""
    safe = Path(filename).name
    path = FILES_DIR / safe
    if not path.exists():
        raise HTTPException(404, "File not found")
    if path.suffix.lower() not in (".xlsx", ".xls"):
        raise HTTPException(400, "Not an Excel file")
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        try:
            for name in wb.sheetnames[:5]:
                ws = wb[name]
                rows_data = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    rows_data.append([str(c) if c is not None else "" for c in row])
                    if i > 500:
                        break
                if rows_data:
                    sheets.append({"name": name, "headers": rows_data[0], "rows": rows_data[1:]})
        finally:
            wb.close()
        return {"filename": safe, "sheets": sheets}
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(500, "Failed to read Excel file")


@app.get("/render-pptx/{filename}")
def render_pptx(filename: str):
    """Convert PPTX to HTML slideshow."""
    import html as _html

    safe = Path(filename).name
    path = FILES_DIR / safe
    if not path.exists():
        raise HTTPException(404, "File not found")
    if path.suffix.lower() != ".pptx":
        raise HTTPException(400, "Not a PPTX file")
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides_html = []
        slides_text = []
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides, 1):
            paras = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            paras.append(t)
            escaped = [f"<p>{_html.escape(p)}</p>" for p in paras]
            slides_html.append(
                f'<div class="slide" id="s{i}">'
                f'<div class="sn">{i} / {total}</div>'
                f'<div class="sc">{"".join(escaped)}</div>'
                f"</div>"
            )
            if paras:
                slides_text.append(f"Slide {i}: " + " | ".join(paras))

        html_out = (
            '<!DOCTYPE html><html><head><meta charset="utf-8">'
            "<style>"
            "*{box-sizing:border-box;margin:0;padding:0}"
            "body{background:#111827;font-family:system-ui,sans-serif;height:100vh;overflow:hidden}"
            ".slide{display:none;width:100%;height:100vh;padding:48px 56px;position:relative;background:#1e293b;color:#e2e8f0;flex-direction:column;justify-content:center}"
            ".slide.active{display:flex}"
            ".sn{position:absolute;top:12px;right:16px;font-size:11px;color:#64748b}"
            ".sc p{margin-bottom:14px;font-size:18px;line-height:1.6}"
            ".sc p:first-child{font-size:28px;font-weight:700;color:#93c5fd;margin-bottom:24px}"
            ".ctrl{position:fixed;bottom:16px;left:50%;transform:translateX(-50%);display:flex;gap:10px;z-index:10}"
            ".ctrl button{background:#1d4ed8;color:#fff;border:none;padding:7px 20px;border-radius:6px;cursor:pointer;font-size:13px}"
            ".ctrl button:hover{background:#2563eb}"
            "</style></head><body>"
            + "".join(slides_html)
            + '<div class="ctrl"><button onclick="p()">&#8592; Prev</button>'
            + '<button onclick="n()">Next &#8594;</button></div>'
            + "<script>"
            + "let c=0;const s=document.querySelectorAll('.slide');s[0].classList.add('active');"
            + "function show(i){s[c].classList.remove('active');c=(i+s.length)%s.length;s[c].classList.add('active');}"
            + "function n(){show(c+1)}function p(){show(c-1)}"
            + "document.addEventListener('keydown',e=>{if(e.key==='ArrowRight'||e.key===' ')n();if(e.key==='ArrowLeft')p();});"
            + "</script></body></html>"
        )
        return {"html": html_out, "text": "\n".join(slides_text), "slides": total}
    except ImportError:
        raise HTTPException(500, "python-pptx not installed")
    except Exception as e:
        raise HTTPException(500, f"PPTX render failed: {e}")


@app.get("/files")
@app.get("/admin/files")
def list_files():
    if not FILES_DIR.exists():
        return {"files": []}
    files = [
        {"name": f.name, "size_kb": f.stat().st_size // 1024, "ext": f.suffix.lstrip(".")}
        for f in sorted(FILES_DIR.iterdir())
        if f.is_file()
    ]
    return {"files": files}


@app.get("/download/{filename}")
def download_file(filename: str):
    safe = Path(filename).name
    path = FILES_DIR / safe
    if not path.exists() or not path.is_file():
        raise HTTPException(404, "file not found")
    return FileResponse(str(path), filename=safe)


@app.get("/admin/memory")
def admin_list_memory(_: None = Depends(_require_trusted_ip)):
    return {"memory": list_memory()}


@app.get("/admin/memory.md")
def admin_memory_markdown(_: None = Depends(_require_trusted_ip)):
    from fastapi.responses import PlainTextResponse
    return PlainTextResponse(get_all_memory(), media_type="text/markdown")


@app.delete("/admin/memory/{key}")
def admin_delete_memory(key: str, _: None = Depends(_require_trusted_ip)):
    ok = delete_memory(key)
    if not ok:
        raise HTTPException(404, "key not found")
    return {"status": "deleted", "key": key}


@app.post("/admin/memory")
def admin_store_memory(body: MemoryRequest, _: None = Depends(_require_trusted_ip)):
    store_memory(body.key.strip(), body.value.strip(), body.category.strip())
    return {"status": "ok", "key": body.key}


@app.get("/conversations")
def conversations():
    return {"conversations": list_recent_conversations()}


@app.get("/history/{conversation_id}")
def history(conversation_id: str):
    if not UUID_RE.match(conversation_id):
        raise HTTPException(400, "invalid conversation_id format")
    return {"messages": get_history(conversation_id)}
